import streamlit as st
import docx
from pptx import Presentation
import fitz  # Correct import for PyMuPDF
import openai
import os

openai.api_key = os.getenv("OPENAI_API_KEY")

def extract_text_from_pdf(uploaded_file):
    text = ""
    with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(uploaded_file):
    doc = docx.Document(uploaded_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def generate_questions(text, num_questions=5):
    prompt = f"Generate {num_questions} quiz questions with answers from the following content:\n\n{text}"
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}]
    )
    questions_with_answers = response['choices'][0]['message']['content']
    # Format questions so answers appear on the next line
    formatted_questions = "\n\n".join(
        [q.replace("Answer:", "\nAnswer:") for q in questions_with_answers.split("\n\n")]
    )
    return formatted_questions

st.title("📚 Document to Quiz Generator")

if not openai.api_key:
    st.error("Missing OPENAI_API_KEY. Set it in your environment before running the app.")
    st.stop()

uploaded_file = st.file_uploader("Upload a document (PDF, PPTX, DOCX)", type=["pdf", "pptx", "docx"])
if uploaded_file:
    if uploaded_file.name.endswith(".pdf"):
        text = extract_text_from_pdf(uploaded_file)
    elif uploaded_file.name.endswith(".docx"):
        text = extract_text_from_docx(uploaded_file)
    elif uploaded_file.name.endswith(".pptx"):
        text = extract_text_from_pptx(uploaded_file)
    else:
        st.error("Unsupported file format")
        st.stop()

    if text:
        st.success("✅ Document processed successfully!")
        num_questions = st.number_input("Enter the number of questions to generate:", min_value=1, step=1, value=5)
        if st.button("Generate Quiz"):
            with st.spinner("Generating questions..."):
                quiz = generate_questions(text, num_questions=num_questions)
                st.markdown("### 📋 Quiz Questions")
                st.markdown(quiz)
